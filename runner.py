#!/usr/bin/env python3
"""
pptx_caract.py — v2.0 (Python puro, sin Node.js)
Genera presentación PowerPoint de caracterización al ingreso (TOP1)
7 slides · Compatible con cualquier país TOP
"""
import glob, os, unicodedata, io, warnings
import pandas as pd
import numpy as np
import matplotlib; matplotlib.use('Agg')
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
warnings.filterwarnings('ignore')

# ── Rutas (inyectadas por runner) ─────────────────────────────────────────────
INPUT_FILE    = None   # runner inyecta la ruta real
OUTPUT_FILE   = None   # runner inyecta la ruta real
SHEET_NAME    = 'Base Wide'
FILTRO_CENTRO = None   # runner inyecta el filtro si aplica
NOMBRE_SERVICIO = 'Servicio de Tratamiento'
PERIODO         = ''

# ── Colores ───────────────────────────────────────────────────────────────────
C_DARK  = RGBColor(0x1F, 0x38, 0x64)
C_MID   = RGBColor(0x2E, 0x75, 0xB6)
C_LIGHT = RGBColor(0xBD, 0xD7, 0xEE)
C_ACC   = RGBColor(0x00, 0xB0, 0xF0)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY  = RGBColor(0x59, 0x59, 0x59)
C_GREEN = RGBColor(0x53, 0x81, 0x35)
PIE_COLS = ['#2E75B6','#1F3864','#4472C4','#9DC3E6','#00B0F0',
            '#538135','#BFBFBF','#C00000','#ED7D31']

SLIDE_W = Inches(10); SLIDE_H = Inches(5.625)

def _norm(s):
    return unicodedata.normalize('NFD', str(s).lower()).encode('ascii','ignore').decode()

_PAISES = {
    'republica_dominicana':'República Dominicana','dominicana':'República Dominicana',
    'honduras':'Honduras','panama':'Panamá','panam':'Panamá',
    'el_salvador':'El Salvador','salvador':'El Salvador',
    'mexico':'México','mexic':'México','ecuador':'Ecuador',
    'peru':'Perú','argentina':'Argentina','colombia':'Colombia',
    'chile':'Chile','bolivia':'Bolivia','paraguay':'Paraguay',
    'uruguay':'Uruguay','venezuela':'Venezuela','guatemala':'Guatemala',
    'costa_rica':'Costa Rica','costarica':'Costa Rica','nicaragua':'Nicaragua',
}
def _extraer_pais(fn):
    f = _norm(str(fn).replace('.','_'))
    for k,v in _PAISES.items():
        if k in f: return v
    return None

def _detectar_pais(wide_file):
    try:
        rs = pd.read_excel(wide_file, sheet_name='Resumen', header=None)
        for _, row in rs.iterrows():
            for v in row.tolist():
                p = _extraer_pais(str(v))
                if p: return p
    except: pass
    return _extraer_pais(os.path.basename(wide_file))

def auto_archivo_wide():
    candidatos = (
        glob.glob('/mnt/user-data/uploads/*Wide*.xlsx') +
        glob.glob('/mnt/user-data/uploads/*wide*.xlsx') +
        glob.glob('/mnt/user-data/uploads/TOP_Base*.xlsx') +
        glob.glob('/home/claude/TOP_Base_Wide.xlsx'))
    if not candidatos: raise FileNotFoundError('No se encontró la base Wide TOP.')
    return candidatos[0]

# ── Helpers PPT ───────────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill, alpha=None):
    from pptx.util import Inches
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape

def add_txt(slide, text, x, y, w, h, size=11, bold=False, color=None,
            align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = str(text)
    run.font.size = Pt(size); run.font.bold = bold; run.font.italic = italic
    if color: run.font.color.rgb = color
    return txb

def header(slide, titulo):
    add_rect(slide, 0, 0, 10, 0.72, C_DARK)
    add_rect(slide, 5.5, 0, 4.5, 0.72, C_MID)
    add_txt(slide, titulo, 0.25, 0.05, 9.5, 0.62,
            size=18, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)

def footer(slide, txt):
    add_txt(slide, txt, 0.25, 5.32, 9.5, 0.25,
            size=8, color=C_GRAY, align=PP_ALIGN.CENTER, italic=True)

def fig_to_pptx(slide, fig, x, y, w, h):
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=130, bbox_inches='tight', facecolor='white')
    buf.seek(0); plt.close(fig)
    slide.shapes.add_picture(buf, Inches(x), Inches(y), Inches(w), Inches(h))

def _ax_style(ax):
    ax.set_facecolor('white'); ax.yaxis.grid(True, color='#E2E8F0', linewidth=0.6, zorder=0)
    ax.set_axisbelow(True)
    for sp in ['top','right']: ax.spines[sp].set_visible(False)
    ax.spines['left'].set_color('#D0D0D0'); ax.spines['bottom'].set_color('#D0D0D0')

# ── Detección de columnas ─────────────────────────────────────────────────────
def detectar_columnas(cols):
    sust_cols = []
    for c in cols:
        if c.endswith('_TOP1') and 'Total (0-28)' in c:
            base = c.replace('_TOP1','')
            if base.startswith('1)'):
                partes = base.split('>>')
                if len(partes) >= 3:
                    nombre = partes[-2].strip().split('(')[0].strip()
                    sust_cols.append((nombre, c))
    tr_sn = []
    for c in cols:
        if c.endswith('_TOP1') and c.replace('_TOP1','').startswith('3)') and '>>' in c:
            nombre = c.replace('_TOP1','').split('>>')[-1].strip()
            tr_sn.append((nombre, c))
    vif     = next((c for c in cols if c.endswith('_TOP1') and '4)' in c
                    and 'Violencia Intrafamiliar' in c and 'Total (0-28)' in c), None)
    sal_psi = next((c for c in cols if c.endswith('_TOP1') and c.replace('_TOP1','').startswith('6)')), None)
    sal_fis = next((c for c in cols if c.endswith('_TOP1') and c.replace('_TOP1','').startswith('8)')), None)
    cal_vid = next((c for c in cols if c.endswith('_TOP1') and c.replace('_TOP1','').startswith('10)')), None)
    viv1    = next((c for c in cols if c.endswith('_TOP1') and '9)' in c and 'estable' in c.lower()), None)
    viv2    = next((c for c in cols if c.endswith('_TOP1') and '9)' in c and 'básicas' in c.lower()), None)
    sust_pp = next((c for c in cols if c.endswith('_TOP1') and c.replace('_TOP1','').startswith('2)')
                    and 'sustancia principal' in c.lower()), None)
    sexo    = next((c for c in cols if c.endswith('_TOP1') and 'sexo' in c.lower()), None)
    fn_col  = next((c for c in cols if c.endswith('_TOP1') and 'nacimiento' in c.lower()), None)
    fecha   = next((c for c in cols if c.endswith('_TOP1') and 'fecha entrevista' in c.lower()), None)
    return dict(sust_cols=sust_cols, tr_sn=tr_sn, vif=vif,
                sal_psi=sal_psi, sal_fis=sal_fis, cal_vid=cal_vid,
                viv1=viv1, viv2=viv2, sust_pp=sust_pp,
                sexo=sexo, fn_col=fn_col, fecha=fecha)

def norm_sust(s):
    if pd.isna(s) or str(s).strip() == '0': return None
    s = str(s).strip().lower()
    if any(x in s for x in ['alcohol','cerveza','licor','aguard']): return 'Alcohol'
    if any(x in s for x in ['marihu','cannabis','marij']):          return 'Marihuana'
    if any(x in s for x in ['pasta base','pasta','papelillo']):     return 'Pasta Base'
    if any(x in s for x in ['crack','cristal','piedra','paco']):    return 'Crack/Cristal'
    if any(x in s for x in ['cocain','perico']):                    return 'Cocaína'
    if any(x in s for x in ['tabaco','cigarr','nicot']):            return 'Tabaco'
    if any(x in s for x in ['inhalant','thiner','activo']):         return 'Inhalantes'
    if any(x in s for x in ['sedant','benzod','tranqui']):          return 'Sedantes'
    if any(x in s for x in ['opiod','heroina','morfin','fentanil']): return 'Opiáceos'
    if any(x in s for x in ['metanfet','anfetam']):                 return 'Metanfetamina'
    return 'Otras'

def _es_positivo(valor):
    s = str(valor).strip().lower()
    if s in ('sí','si'): return True
    if s in ('no','no aplica','nunca','nan',''): return False
    n = pd.to_numeric(valor, errors='coerce')
    return not pd.isna(n) and n > 0

# ── Carga y cálculo de datos ──────────────────────────────────────────────────
def cargar_datos():
    global NOMBRE_SERVICIO, PERIODO, INPUT_FILE, OUTPUT_FILE, FILTRO_CENTRO
    # Leer de variables de entorno si fueron inyectadas por el runner
    if INPUT_FILE is None:
        INPUT_FILE = os.environ.get('QALAT_WIDE') or auto_archivo_wide()
    if OUTPUT_FILE is None:
        OUTPUT_FILE = os.environ.get('QALAT_OUT', '/home/claude/TOP_Presentacion_Caracterizacion.pptx')
    if FILTRO_CENTRO is None:
        FILTRO_CENTRO = os.environ.get('QALAT_CENTRO') or None
    _pais = _detectar_pais(INPUT_FILE)
    NOMBRE_SERVICIO = _pais if _pais else 'Servicio de Tratamiento'

    _periodo_auto = None
    try:
        _rs = pd.read_excel(INPUT_FILE, sheet_name='Resumen', header=None)
        for _, _row in _rs.iterrows():
            for _v in _row.tolist():
                if 'Período' in str(_v) or 'periodo' in str(_v).lower(): continue
                if '–' in str(_v) or (' ' in str(_v) and any(
                        m in str(_v) for m in ['Enero','Feb','Mar','Abr','May','Jun',
                                               'Jul','Ago','Sep','Oct','Nov','Dic','2024','2025','2026'])):
                    _periodo_auto = str(_v).strip(); break
            if _periodo_auto: break
    except: pass
    PERIODO = _periodo_auto if _periodo_auto else '2025'

    df = pd.read_excel(INPUT_FILE, sheet_name=SHEET_NAME, header=1)
    df.columns = [str(c) for c in df.columns]
    cols = df.columns.tolist()

    _col_centro = next((c for c in cols if any(x in _norm(c) for x in
                        ['codigo del centro','servicio de tratamiento',
                         'centro/ servicio','codigo centro'])), None)
    if FILTRO_CENTRO and _col_centro:
        df = df[df[_col_centro].astype(str).str.strip()==FILTRO_CENTRO].copy().reset_index(drop=True)
        _pl = _detectar_pais(INPUT_FILE)
        NOMBRE_SERVICIO = f'{_pl}  —  Centro {FILTRO_CENTRO}' if _pl else f'Centro {FILTRO_CENTRO}'

    N  = len(df)
    DC = detectar_columnas(cols)
    hoy = pd.Timestamp.now()

    # Sexo
    n_h = n_m = 0; pct_h = pct_m = 0.0
    if DC['sexo']:
        sc = df[DC['sexo']].astype(str).str.strip().str.upper()
        nv = int(sc.isin(['H','M']).sum())
        n_h = int((sc=='H').sum()); n_m = int((sc=='M').sum())
        pct_h = round(n_h/nv*100,1) if nv else 0
        pct_m = round(n_m/nv*100,1) if nv else 0

    # Edad
    edad_media = 0; edad_grupos = []
    if DC['fn_col'] and DC['fecha']:
        fn  = pd.to_datetime(df[DC['fn_col']], errors='coerce')
        ref = pd.to_datetime(df[DC['fecha']], errors='coerce').fillna(hoy)
        edad = ((ref-fn).dt.days/365.25).round(1)
        edad = edad[(edad>=10)&(edad<=100)]
        if len(edad):
            edad_media = round(float(edad.mean()),1)
            bins = [0,17,30,40,50,60,200]; labs = ['<18','18–30','31–40','41–50','51–60','61+']
            ec = pd.cut(edad, bins=bins, labels=labs)
            total_e = len(edad)
            edad_grupos = [{'label':l,'n':int((ec==l).sum()),
                            'pct':round(int((ec==l).sum())/total_e*100,1)}
                           for l in labs if int((ec==l).sum())>0]

    # Sustancia principal
    sust_ppal = []; sust_top1 = '—'; sust_top1_pct = 0
    if DC['sust_pp']:
        sr = df[DC['sust_pp']].apply(norm_sust).dropna()
        vc = sr.value_counts(); total_sp = len(sr)
        sust_ppal = [{'label':k,'pct':round(v/total_sp*100,1),'n':int(v)} for k,v in vc.items()]
        sust_top1 = vc.index[0] if len(vc) else '—'
        sust_top1_pct = round(vc.iloc[0]/total_sp*100,1) if len(vc) else 0

    # Días consumo por sustancia principal
    sust_norm = df[DC['sust_pp']].apply(norm_sust) if DC['sust_pp'] else pd.Series([None]*N)
    dias_princ = []
    for lbl, col in DC['sust_cols']:
        v = pd.to_numeric(df[col], errors='coerce')
        mask = sust_norm.apply(lambda s: isinstance(s,str) and lbl.lower() in s.lower())
        sub  = v[mask & (v>0)].dropna()
        if len(sub): dias_princ.append({'label':lbl,'prom':round(float(sub.mean()),1),'n':int(len(sub))})
    dias_princ.sort(key=lambda x:-x['prom'])

    # % Consumidores
    consumo_pct = []
    for lbl, col in DC['sust_cols']:
        v = pd.to_numeric(df[col], errors='coerce').fillna(0)
        n_c = int((v>0).sum())
        if n_c > 0: consumo_pct.append({'label':lbl,'pct':round(n_c/N*100,1),'n':n_c})
    consumo_pct.sort(key=lambda x:-x['pct'])

    # Días por sustancia
    dias_sust = []
    for lbl, col in DC['sust_cols']:
        v = pd.to_numeric(df[col], errors='coerce'); sub = v[v>0].dropna()
        if len(sub): dias_sust.append({'label':lbl,'prom':round(float(sub.mean()),1),'n':int(len(sub))})
    dias_sust.sort(key=lambda x:-x['prom'])

    # Transgresión
    def has_tr(row):
        for _, c in DC['tr_sn']:
            if _es_positivo(row.get(c,'')): return True
        if DC['vif']:
            v = pd.to_numeric(row.get(DC['vif'], np.nan), errors='coerce')
            return not np.isnan(v) and v > 0
        return False
    t = df.apply(lambda r: int(has_tr(r)), axis=1)
    n_tr = int(t.sum()); pct_tr = round(n_tr/N*100,1)
    tipos_tr = []
    for lbl, col in DC['tr_sn']:
        n = int(df[col].apply(_es_positivo).sum())
        tipos_tr.append({'label':lbl,'n':n,'pct':round(n/N*100,1)})
    if DC['vif']:
        vif_v = pd.to_numeric(df[DC['vif']], errors='coerce')
        n_vif = int((vif_v>0).sum())
        tipos_tr.append({'label':'VIF','n':n_vif,'pct':round(n_vif/N*100,1)})
    tipos_tr = [t for t in tipos_tr if t['pct']>0]

    # Salud
    salud = []
    for lbl, col in [('Salud Psicológica',DC['sal_psi']),
                     ('Salud Física',DC['sal_fis']),
                     ('Calidad de Vida',DC['cal_vid'])]:
        if col:
            v = pd.to_numeric(df[col], errors='coerce')
            salud.append({'label':lbl,'prom':round(float(v.mean()),1),'nv':int(v.notna().sum())})

    # Vivienda
    def viv(col):
        if not col: return {'n':0,'pct':0}
        nv_ = int(df[col].isin(['Sí','No']).sum()) or N
        n_  = int((df[col]=='Sí').sum())
        return {'n':n_,'pct':round(n_/nv_*100,1)}
    viv1 = viv(DC['viv1']); viv2 = viv(DC['viv2'])

    return dict(N=N, n_h=n_h, n_m=n_m, pct_h=pct_h, pct_m=pct_m,
                edad_media=edad_media, edad_grupos=edad_grupos,
                sust_ppal=sust_ppal, sust_top1=sust_top1, sust_top1_pct=sust_top1_pct,
                dias_princ=dias_princ, consumo_pct=consumo_pct, dias_sust=dias_sust,
                n_tr=n_tr, pct_tr=pct_tr, tipos_tr=tipos_tr,
                salud=salud, viv1=viv1, viv2=viv2)

# ── Gráficos ──────────────────────────────────────────────────────────────────
def g_sexo(d):
    fig, ax = plt.subplots(figsize=(4,3))
    vals = [d['n_h'], d['n_m']]; labs = ['Hombre','Mujer']
    cols = ['#2E75B6','#00B0F0']
    bars = ax.bar(labs, vals, color=cols, width=0.5, zorder=3)
    for bar, v, p in zip(bars, vals, [d['pct_h'],d['pct_m']]):
        ax.text(bar.get_x()+bar.get_width()/2, bar.get_height()+0.3,
                f'{v}\n({p}%)', ha='center', va='bottom', fontsize=10, fontweight='bold', color='#333')
    ax.set_ylim(0, max(vals)*1.35 if max(vals)>0 else 1)
    ax.set_ylabel('N personas', fontsize=8, color='#595959')
    _ax_style(ax); fig.patch.set_facecolor('white'); fig.tight_layout()
    return fig

def g_edad(d):
    if not d['edad_grupos']: return None
    labs = [g['label'] for g in d['edad_grupos']]
    vals = [g['n'] for g in d['edad_grupos']]
    fig, ax = plt.subplots(figsize=(4.5,3))
    cols = ['#2E75B6' if v==max(vals) else '#BDD7EE' for v in vals]
    bars = ax.barh(labs, vals, color=cols, zorder=3)
    total = sum(vals)
    for bar, v in zip(bars, vals):
        pct = round(v/total*100,1) if total else 0
        ax.text(bar.get_width()+0.2, bar.get_y()+bar.get_height()/2,
                f'{v} ({pct}%)', va='center', fontsize=8, color='#333')
    ax.set_xlim(0, max(vals)*1.5 if max(vals)>0 else 1)
    _ax_style(ax); fig.patch.set_facecolor('white'); fig.tight_layout()
    return fig

def g_torta(d):
    if not d['sust_ppal']: return None
    labs = [s['label'] for s in d['sust_ppal']]
    vals = [s['n'] for s in d['sust_ppal']]
    fig, ax = plt.subplots(figsize=(5,4))
    wedges, _, autotexts = ax.pie(vals, labels=None, colors=PIE_COLS[:len(vals)],
        autopct=lambda p: f'{p:.1f}%' if p>3 else '', startangle=140, pctdistance=0.72,
        wedgeprops={'edgecolor':'white','linewidth':2})
    for at in autotexts: at.set_fontsize(9); at.set_color('white'); at.set_fontweight('bold')
    ax.legend(wedges, [f'{l} (n={v})' for l,v in zip(labs,vals)],
              loc='lower center', bbox_to_anchor=(0.5,-0.18), ncol=2, fontsize=8, frameon=False)
    ax.set_aspect('equal'); fig.patch.set_facecolor('white'); fig.tight_layout()
    return fig

def g_barras_h(datos, ylabel=''):
    if not datos: return None
    labs = [d['label'] for d in datos]; vals = [d['prom'] for d in datos]
    fig, ax = plt.subplots(figsize=(max(4, len(labs)*0.9), 3))
    cols = ['#2E75B6' if v==max(vals) else '#BDD7EE' for v in vals]
    bars = ax.bar(labs, vals, color=cols, width=0.55, zorder=3)
    for bar, d_item in zip(bars, datos):
        ax.text(bar.get_x()+bar.get_width()/2, bar.get_height()+0.2,
                f'{d_item["prom"]}d\n(n={d_item["n"]})', ha='center', va='bottom',
                fontsize=8, fontweight='bold', color='#333')
    ax.set_ylim(0, max(vals)*1.4 if max(vals)>0 else 1)
    ax.set_ylabel(ylabel, fontsize=8, color='#595959')
    ax.tick_params(labelsize=8)
    _ax_style(ax); fig.patch.set_facecolor('white'); fig.tight_layout()
    return fig

def g_consumo(d):
    if not d['consumo_pct']: return None
    labs = [x['label'] for x in d['consumo_pct']]
    vals = [x['pct'] for x in d['consumo_pct']]
    fig, ax = plt.subplots(figsize=(max(4, len(labs)*0.9), 3))
    cols = ['#2E75B6' if v==max(vals) else '#BDD7EE' for v in vals]
    bars = ax.bar(labs, vals, color=cols, width=0.55, zorder=3)
    for bar, item in zip(bars, d['consumo_pct']):
        ax.text(bar.get_x()+bar.get_width()/2, bar.get_height()+0.5,
                f'{item["pct"]}%\n(n={item["n"]})', ha='center', va='bottom',
                fontsize=8, fontweight='bold', color='#333')
    ax.set_ylim(0, max(vals)*1.4 if max(vals)>0 else 1)
    ax.set_ylabel('% de personas', fontsize=8, color='#595959')
    ax.tick_params(labelsize=8)
    _ax_style(ax); fig.patch.set_facecolor('white'); fig.tight_layout()
    return fig

def g_transgresion(d):
    fig, ax = plt.subplots(figsize=(3.5,3))
    n_no = d['N'] - d['n_tr']
    pct_no = round(n_no/d['N']*100,1) if d['N'] else 0
    bars = ax.bar(['Con\ntransgresión','Sin\ntransgresión'],
                  [d['pct_tr'], pct_no], color=['#2E75B6','#BDD7EE'], width=0.5, zorder=3)
    for bar, v, n in zip(bars, [d['pct_tr'], pct_no], [d['n_tr'], n_no]):
        ax.text(bar.get_x()+bar.get_width()/2, bar.get_height()+0.5,
                f'{n}\n({v}%)', ha='center', va='bottom', fontsize=10, fontweight='bold', color='#333')
    ax.set_ylim(0, 120); ax.set_ylabel('% personas', fontsize=8, color='#595959')
    _ax_style(ax); fig.patch.set_facecolor('white'); fig.tight_layout()
    return fig

def g_tipos_tr(d):
    if not d['tipos_tr']: return None
    labs = [t['label'] for t in d['tipos_tr']]; vals = [t['pct'] for t in d['tipos_tr']]
    fig, ax = plt.subplots(figsize=(max(4,len(labs)*0.9), 3))
    cols = ['#2E75B6' if v==max(vals) else '#BDD7EE' for v in vals]
    bars = ax.bar(labs, vals, color=cols, width=0.55, zorder=3)
    for bar, item in zip(bars, d['tipos_tr']):
        ax.text(bar.get_x()+bar.get_width()/2, bar.get_height()+0.4,
                f'{item["pct"]}%', ha='center', va='bottom', fontsize=9, fontweight='bold', color='#333')
    ax.set_ylim(0, max(vals)*1.4 if max(vals)>0 else 1)
    ax.set_ylabel('% personas', fontsize=8, color='#595959')
    ax.tick_params(labelsize=8)
    _ax_style(ax); fig.patch.set_facecolor('white'); fig.tight_layout()
    return fig

def g_salud(d):
    if not d['salud']: return None
    labs = [s['label'] for s in d['salud']]; vals = [s['prom'] for s in d['salud']]
    fig, ax = plt.subplots(figsize=(4, 2.5))
    bars = ax.barh(labs, vals, color='#2E75B6', height=0.45, zorder=3)
    for bar, v in zip(bars, vals):
        ax.text(bar.get_width()+0.2, bar.get_y()+bar.get_height()/2,
                f'{v}/20', va='center', fontsize=9, fontweight='bold', color='#333')
    ax.set_xlim(0, 22); ax.axvline(x=10, color='#BFBFBF', linestyle='--', linewidth=0.8)
    _ax_style(ax); fig.patch.set_facecolor('white'); fig.tight_layout()
    return fig

def g_vivienda(d):
    fig, ax = plt.subplots(figsize=(3.5, 2.5))
    cats = ['Lugar\nestable','Condiciones\nbásicas']
    vals = [d['viv1']['pct'], d['viv2']['pct']]
    bars = ax.barh(cats, vals, color='#2E75B6', height=0.45, zorder=3)
    for bar, v in zip(bars, vals):
        ax.text(bar.get_width()+1, bar.get_y()+bar.get_height()/2,
                f'{v}%', va='center', fontsize=10, fontweight='bold', color='#333')
    ax.set_xlim(0, 115)
    _ax_style(ax); fig.patch.set_facecolor('white'); fig.tight_layout()
    return fig

# ── Construcción del PPT ──────────────────────────────────────────────────────
def build_pptx(d):
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]  # layout en blanco

    titulo_base = f'Caracterización al Ingreso · {NOMBRE_SERVICIO}'
    pie_txt = f'N = {d["N"]}  ·  {NOMBRE_SERVICIO}  ·  {PERIODO}'

    # ── SLIDE 1: PORTADA ──────────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    add_rect(sl, 0, 0, 3.8, 5.625, C_DARK)
    add_rect(sl, 3.0, 0, 1.5, 5.625, C_MID)
    add_txt(sl, 'Caracterización', 0.25, 1.5, 3.2, 0.7, size=22, bold=True, color=C_WHITE)
    add_txt(sl, 'Ingreso a Tratamiento · TOP', 0.25, 2.3, 3.2, 0.5, size=12, color=C_LIGHT)
    add_txt(sl, NOMBRE_SERVICIO.upper(), 0.25, 3.1, 3.2, 0.6, size=13, bold=True, color=C_WHITE)
    add_txt(sl, PERIODO, 0.25, 3.75, 3.2, 0.4, size=11, color=C_LIGHT)
    add_txt(sl, f'N = {d["N"]} personas al ingreso', 0.25, 4.3, 3.2, 0.4, size=10, color=C_LIGHT)
    add_txt(sl, 'Monitoreo de Resultados\nde Tratamiento', 4.5, 2.2, 5.2, 1.0,
            size=16, bold=True, color=C_DARK, align=PP_ALIGN.LEFT)
    add_txt(sl, f'Sustancia principal: {d["sust_top1"]} ({d["sust_top1_pct"]}%)',
            4.5, 3.3, 5.2, 0.4, size=11, color=C_MID)
    add_txt(sl, f'{d["pct_h"]}% Hombres  ·  Edad promedio: {d["edad_media"]} años',
            4.5, 3.75, 5.2, 0.4, size=10, color=C_GRAY)

    # ── SLIDE 2: ANTECEDENTES ─────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    header(sl, titulo_base)
    # Línea divisora
    line = sl.shapes.add_shape(1, Inches(4.95), Inches(0.78), Inches(0.02), Inches(4.85))
    line.fill.solid(); line.fill.fore_color.rgb = RGBColor(0xD9,0xD9,0xD9)
    line.line.fill.background()

    add_txt(sl, 'DISTRIBUCIÓN POR SEXO', 0.25, 0.82, 4.5, 0.35,
            size=11, bold=True, color=C_MID)
    fig_to_pptx(sl, g_sexo(d), 0.25, 1.2, 4.5, 3.8)

    add_txt(sl, 'DISTRIBUCIÓN POR EDAD', 5.15, 0.82, 4.6, 0.35,
            size=11, bold=True, color=C_MID)
    fig_edad = g_edad(d)
    if fig_edad: fig_to_pptx(sl, fig_edad, 5.1, 1.2, 4.7, 3.8)
    footer(sl, pie_txt)

    # ── SLIDE 3: SUSTANCIA PRINCIPAL ──────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    header(sl, titulo_base)
    add_txt(sl, f'SUSTANCIA PRINCIPAL AL INGRESO  ·  {d["sust_top1"]} ({d["sust_top1_pct"]}%)',
            0.25, 0.82, 9.5, 0.35, size=12, bold=True, color=C_MID)
    fig_t = g_torta(d)
    if fig_t: fig_to_pptx(sl, fig_t, 2.0, 1.2, 6.0, 4.0)
    footer(sl, pie_txt)

    # ── SLIDE 4: DÍAS DE CONSUMO SUSTANCIA PRINCIPAL ──────────────────────────
    sl = prs.slides.add_slide(blank)
    header(sl, titulo_base)
    add_txt(sl, 'PROMEDIO DE DÍAS DE CONSUMO · Sustancia principal (últimas 4 semanas)',
            0.25, 0.82, 9.5, 0.35, size=11, bold=True, color=C_MID)
    fig_dp = g_barras_h(d['dias_princ'], ylabel='Promedio días (0–28)')
    if fig_dp: fig_to_pptx(sl, fig_dp, 0.5, 1.2, 9.0, 4.0)
    footer(sl, pie_txt)

    # ── SLIDE 5: % CONSUMIDORES + DÍAS POR SUSTANCIA ──────────────────────────
    sl = prs.slides.add_slide(blank)
    header(sl, titulo_base)
    line = sl.shapes.add_shape(1, Inches(4.95), Inches(0.78), Inches(0.02), Inches(4.85))
    line.fill.solid(); line.fill.fore_color.rgb = RGBColor(0xD9,0xD9,0xD9)
    line.line.fill.background()

    add_txt(sl, '% DE PERSONAS QUE CONSUME\nCada sustancia al ingreso',
            0.25, 0.82, 4.5, 0.55, size=11, bold=True, color=C_MID)
    fig_c = g_consumo(d)
    if fig_c: fig_to_pptx(sl, fig_c, 0.2, 1.45, 4.5, 3.85)

    add_txt(sl, 'PROMEDIO DE DÍAS DE CONSUMO\nPor sustancia (solo consumidores)',
            5.15, 0.82, 4.6, 0.55, size=11, bold=True, color=C_MID)
    fig_ds = g_barras_h(d['dias_sust'], ylabel='Promedio días (0–28)')
    if fig_ds: fig_to_pptx(sl, fig_ds, 5.15, 1.45, 4.6, 3.85)
    footer(sl, pie_txt)

    # ── SLIDE 6: TRANSGRESIÓN ─────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    header(sl, titulo_base)
    line = sl.shapes.add_shape(1, Inches(4.95), Inches(0.78), Inches(0.02), Inches(4.85))
    line.fill.solid(); line.fill.fore_color.rgb = RGBColor(0xD9,0xD9,0xD9)
    line.line.fill.background()

    add_txt(sl, 'Personas que cometieron alguna\ntransgresión a la norma social',
            0.25, 0.82, 4.5, 0.6, size=13, bold=True, color=C_GRAY)
    fig_to_pptx(sl, g_transgresion(d), 0.3, 1.5, 4.4, 3.6)
    add_txt(sl, f'{d["n_tr"]} personas ({d["pct_tr"]}%)',
            0.3, 1.2, 4.4, 0.28, size=12, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)

    add_txt(sl, 'Distribución por tipo de transgresión',
            5.15, 0.82, 4.6, 0.6, size=13, bold=True, color=C_GRAY)
    fig_tt = g_tipos_tr(d)
    if fig_tt: fig_to_pptx(sl, fig_tt, 5.1, 1.5, 4.65, 3.6)
    footer(sl, pie_txt)

    # ── SLIDE 7: SALUD Y VIVIENDA ─────────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    header(sl, titulo_base)
    line = sl.shapes.add_shape(1, Inches(5.05), Inches(0.78), Inches(0.02), Inches(4.85))
    line.fill.solid(); line.fill.fore_color.rgb = RGBColor(0xD9,0xD9,0xD9)
    line.line.fill.background()

    add_txt(sl, 'AUTOPERCEPCIÓN DEL ESTADO DE SALUD\nY CALIDAD DE VIDA (escala 0–20)',
            0.25, 0.82, 4.7, 0.6, size=11, bold=True, color=C_MID)
    fig_sal = g_salud(d)
    if fig_sal: fig_to_pptx(sl, fig_sal, 0.2, 1.5, 4.7, 3.8)

    add_txt(sl, 'CONDICIONES DE VIVIENDA\n(% con condición Sí)',
            5.3, 0.82, 4.5, 0.6, size=11, bold=True, color=C_MID)
    fig_viv = g_vivienda(d)
    if fig_viv: fig_to_pptx(sl, fig_viv, 5.2, 1.5, 4.5, 3.8)
    footer(sl, pie_txt)

    prs.save(OUTPUT_FILE)
    print(f'  ✓ PPT generado: {OUTPUT_FILE}')

# ══════════════════════════════════════════════════════════════════════════════
if __name__ == '__main__':
    print('='*60)
    print('  PPTX Caracterización TOP  v2.0  —  Iniciando...')
    print('='*60)
    d = cargar_datos()
    print(f'  N={d["N"]} | {d["sust_top1"]} {d["sust_top1_pct"]}% | {NOMBRE_SERVICIO}')
    build_pptx(d)
    print(f'\n{"="*60}')
    print(f'  ✅  LISTO  →  {OUTPUT_FILE}')
    print(f'{"="*60}')
