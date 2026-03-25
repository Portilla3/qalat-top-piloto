#!/usr/bin/env python3
"""
pptx_seg.py — v2.0 (Python puro, sin Node.js)
Genera presentación PowerPoint de seguimiento TOP1 vs TOP2
6 slides · Compatible con cualquier país TOP
"""
import glob, os, unicodedata, io, warnings
import pandas as pd
import numpy as np
import matplotlib; matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
warnings.filterwarnings('ignore')

# ── Rutas (inyectadas por runner) ─────────────────────────────────────────────
INPUT_FILE    = None
OUTPUT_FILE   = None
SHEET_NAME    = 'Base Wide'
FILTRO_CENTRO = None
NOMBRE_SERVICIO = 'Servicio de Tratamiento'
PERIODO         = ''

# ── Colores ───────────────────────────────────────────────────────────────────
C_DARK  = RGBColor(0x1F, 0x38, 0x64)
C_MID   = RGBColor(0x2E, 0x75, 0xB6)
C_LIGHT = RGBColor(0xBD, 0xD7, 0xEE)
C_ACC   = RGBColor(0x00, 0xB0, 0xF0)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY  = RGBColor(0x59, 0x59, 0x59)
MC_T1 = '#2E75B6'; MC_T2 = '#00B0F0'
MC_ABS = '#1F3864'; MC_DIS = '#2E75B6'; MC_SC = '#9DC3E6'; MC_EMP = '#BDD7EE'
PIE_COLS = ['#2E75B6','#1F3864','#4472C4','#9DC3E6','#00B0F0','#538135','#BFBFBF','#C00000','#ED7D31']

SLIDE_W = Inches(10); SLIDE_H = Inches(5.625)

def _norm(s):
    return unicodedata.normalize('NFD', str(s).lower()).encode('ascii','ignore').decode()

_PAISES = {
    'republica_dominicana':'República Dominicana','dominicana':'República Dominicana',
    'honduras':'Honduras','panama':'Panamá','panam':'Panamá',
    'el_salvador':'El Salvador','salvador':'El Salvador',
    'mexico':'México','mexic':'México','ecuador':'Ecuador',
    'peru':'Perú','argentina':'Argentina','colombia':'Colombia',
    'chile':'Chile','bolivia':'Bolivia','paraguay':'Paraguay',
    'uruguay':'Uruguay','venezuela':'Venezuela','guatemala':'Guatemala',
    'costa_rica':'Costa Rica','costarica':'Costa Rica','nicaragua':'Nicaragua',
}
def _extraer_pais(fn):
    f = _norm(str(fn).replace('.','_'))
    for k,v in _PAISES.items():
        if k in f: return v
    return None

def _detectar_pais(wide_file):
    try:
        rs = pd.read_excel(wide_file, sheet_name='Resumen', header=None)
        for _, row in rs.iterrows():
            for v in row.tolist():
                p = _extraer_pais(str(v))
                if p: return p
    except: pass
    return _extraer_pais(os.path.basename(wide_file))

def auto_archivo_wide():
    candidatos = (
        glob.glob('/mnt/user-data/uploads/*Wide*.xlsx') +
        glob.glob('/mnt/user-data/uploads/*wide*.xlsx') +
        glob.glob('/mnt/user-data/uploads/TOP_Base*.xlsx') +
        glob.glob('/home/claude/TOP_Base_Wide.xlsx'))
    if not candidatos: raise FileNotFoundError('No se encontró la base Wide TOP.')
    return candidatos[0]

def _es_positivo(valor):
    s = str(valor).strip().lower()
    if s in ('sí','si'): return True
    if s in ('no','no aplica','nunca','nan',''): return False
    n = pd.to_numeric(valor, errors='coerce')
    return not pd.isna(n) and n > 0

def norm_sust(s):
    if pd.isna(s) or str(s).strip() == '0': return None
    s = str(s).strip().lower()
    if any(x in s for x in ['alcohol','cerveza','licor','aguard','alchol','bebida']): return 'Alcohol'
    if any(x in s for x in ['marihu','marjhu','marhuana','cannabis','cannbis']): return 'Cannabis/Marihuana'
    if any(x in s for x in ['pasta base','pasta','papelillo']): return 'Pasta Base'
    if any(x in s for x in ['crack','cristal','piedra','paco']): return 'Crack/Cristal'
    if any(x in s for x in ['cocain','perico','coca ']): return 'Cocaína'
    if any(x in s for x in ['tabaco','cigarr','nicot']): return 'Tabaco'
    if any(x in s for x in ['sedant','benzod','tranqui']): return 'Sedantes'
    if any(x in s for x in ['opiod','heroina','morfin','fentanil']): return 'Opiáceos'
    if any(x in s for x in ['metanfet','anfetam']): return 'Metanfetamina'
    return 'Otras'

# ── Helpers PPT ───────────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape

def add_txt(slide, text, x, y, w, h, size=11, bold=False, color=None,
            align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = str(text)
    run.font.size = Pt(size); run.font.bold = bold; run.font.italic = italic
    if color: run.font.color.rgb = color
    return txb

def header(slide, titulo):
    add_rect(slide, 0, 0, 10, 0.72, C_DARK)
    add_rect(slide, 5.5, 0, 4.5, 0.72, C_MID)
    add_txt(slide, titulo, 0.25, 0.05, 9.5, 0.62, size=18, bold=True, color=C_WHITE)

def footer(slide, txt):
    add_txt(slide, txt, 0.25, 5.32, 9.5, 0.25, size=8, color=C_GRAY,
            align=PP_ALIGN.CENTER, italic=True)

def fig_to_pptx(slide, fig, x, y, w, h):
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=130, bbox_inches='tight', facecolor='white')
    buf.seek(0); plt.close(fig)
    slide.shapes.add_picture(buf, Inches(x), Inches(y), Inches(w), Inches(h))

def _ax_style(ax, horiz=False):
    ax.set_facecolor('white')
    (ax.xaxis if horiz else ax.yaxis).grid(True, color='#E2E8F0', linewidth=0.6, zorder=0)
    ax.set_axisbelow(True)
    for sp in ['top','right']: ax.spines[sp].set_visible(False)
    ax.spines['left'].set_color('#D0D0D0'); ax.spines['bottom'].set_color('#D0D0D0')

def div_v(slide, x):
    line = slide.shapes.add_shape(1, Inches(x), Inches(0.78), Inches(0.02), Inches(4.85))
    line.fill.solid(); line.fill.fore_color.rgb = RGBColor(0xD9,0xD9,0xD9)
    line.line.fill.background()

# ── Detección de columnas ─────────────────────────────────────────────────────
def detectar_columnas(cols):
    col_set = set(cols)
    def par(c1):
        if not c1: return (None, None)
        c2 = c1.replace('_TOP1','_TOP2')
        return (c1, c2 if c2 in col_set else None)

    sust_cols = []
    for c in cols:
        if c.endswith('_TOP1') and 'Total (0-28)' in c:
            base = c.replace('_TOP1','')
            if base.startswith('1)'):
                partes = base.split('>>')
                if len(partes) >= 3:
                    nombre = partes[-2].strip().split('(')[0].strip()
                    c1, c2 = par(c)
                    sust_cols.append((nombre, c1, c2))
    tr_sn = []
    for c in cols:
        if c.endswith('_TOP1') and c.replace('_TOP1','').startswith('3)') and '>>' in c:
            nombre = c.replace('_TOP1','').split('>>')[-1].strip()
            c1, c2 = par(c)
            tr_sn.append((nombre, c1, c2))

    def find(kws):
        for c in cols:
            if not c.endswith('_TOP1'): continue
            nc = _norm(c)
            if all(k in nc for k in kws): return par(c)
        return (None, None)

    vif     = find(['4)','violencia','intrafamiliar','total'])
    sal_psi = find(['6)'])
    sal_fis = find(['8)'])
    cal_vid = find(['10)'])
    viv1    = find(['9)','estable'])
    viv2    = find(['9)','basica'])
    sust_pp = find(['2)','sustancia','principal'])

    return dict(sust_cols=sust_cols, tr_sn=tr_sn, vif=vif,
                sal_psi=sal_psi, sal_fis=sal_fis, cal_vid=cal_vid,
                viv1=viv1, viv2=viv2, sust_pp=sust_pp)

# ── Carga y cálculo de datos ──────────────────────────────────────────────────
def cargar_datos():
    global NOMBRE_SERVICIO, PERIODO, INPUT_FILE, OUTPUT_FILE, FILTRO_CENTRO
    # Leer de variables de entorno si fueron inyectadas por el runner
    if INPUT_FILE is None:
        INPUT_FILE = os.environ.get('QALAT_WIDE') or auto_archivo_wide()
    if OUTPUT_FILE is None:
        OUTPUT_FILE = os.environ.get('QALAT_OUT', '/home/claude/TOP_Presentacion_Seguimiento.pptx')
    if FILTRO_CENTRO is None:
        FILTRO_CENTRO = os.environ.get('QALAT_CENTRO') or None
    _pais = _detectar_pais(INPUT_FILE)
    NOMBRE_SERVICIO = _pais if _pais else 'Servicio de Tratamiento'

    _periodo_auto = None
    try:
        _rs = pd.read_excel(INPUT_FILE, sheet_name='Resumen', header=None)
        for _, _row in _rs.iterrows():
            for _v in _row.tolist():
                if 'Período' in str(_v) or 'periodo' in str(_v).lower(): continue
                if '–' in str(_v) or (' ' in str(_v) and any(
                        m in str(_v) for m in ['Enero','Feb','Mar','Abr','May','Jun',
                                               'Jul','Ago','Sep','Oct','Nov','Dic','2024','2025','2026'])):
                    _periodo_auto = str(_v).strip(); break
            if _periodo_auto: break
    except: pass
    PERIODO = _periodo_auto if _periodo_auto else '2025'

    df = pd.read_excel(INPUT_FILE, sheet_name=SHEET_NAME, header=1)
    df.columns = [str(c) for c in df.columns]

    _col_centro = next((c for c in df.columns if any(x in _norm(c) for x in
                        ['codigo del centro','servicio de tratamiento',
                         'centro/ servicio','codigo centro'])), None)
    if FILTRO_CENTRO and _col_centro:
        df = df[df[_col_centro].astype(str).str.strip()==FILTRO_CENTRO].copy().reset_index(drop=True)
        _pl = _detectar_pais(INPUT_FILE)
        NOMBRE_SERVICIO = f'{_pl}  —  Centro {FILTRO_CENTRO}' if _pl else f'Centro {FILTRO_CENTRO}'

    N_total = len(df)
    seg = df[df['Tiene_TOP2']=='Sí'].copy().reset_index(drop=True)
    N = len(seg)
    DC = detectar_columnas(seg.columns.tolist())

    def pct(n, d): return round(n/d*100,1) if d>0 else 0
    def smean(col):
        if not col or col not in seg.columns: return 0
        v = pd.to_numeric(seg[col], errors='coerce')
        return round(float(v.mean()),1) if v.notna().sum()>0 else 0
    def viv_pct(col):
        if not col or col not in seg.columns: return 0
        nv = int(seg[col].isin(['Sí','No']).sum()) or N
        return pct(int((seg[col]=='Sí').sum()), nv)

    # Tiempo seguimiento
    _fc1 = next((c for c in seg.columns if 'fecha entrevista' in c.lower() and c.endswith('_TOP1')), None)
    _fc2 = next((c for c in seg.columns if 'fecha entrevista' in c.lower() and c.endswith('_TOP2')), None)
    seg_tiempo = {'mediana':None,'min':None,'max':None,'n':0}
    if _fc1 and _fc2:
        _d1 = pd.to_datetime(seg[_fc1], errors='coerce')
        _d2 = pd.to_datetime(seg[_fc2], errors='coerce')
        _dias = (_d2-_d1).dt.days
        _dias_ok = _dias[(_dias>=0)&(_dias<=730)].dropna()
        if len(_dias_ok)>0:
            _m = _dias_ok/30.44
            seg_tiempo = {'mediana':round(float(_m.median()),1),
                          'min':round(float(_m.min()),1),
                          'max':round(float(_m.max()),1),'n':len(_dias_ok)}

    # Sustancia principal
    c1_sp, _ = DC['sust_pp']
    sust = []; sust_top = '—'
    if c1_sp:
        sr1 = seg[c1_sp].apply(norm_sust).dropna()
        nv = len(sr1); vc = sr1.value_counts()
        sust = [{'label':k,'pct':round(v/nv*100,1),'n':int(v)} for k,v in vc.items()]
        sust_top = sust[0]['label'] if sust else '—'

    # Días consumo TOP1 vs TOP2
    dias = []
    for lbl, c1, c2 in DC['sust_cols']:
        v1 = pd.to_numeric(seg[c1], errors='coerce')
        v2 = pd.to_numeric(seg[c2], errors='coerce') if c2 else pd.Series([np.nan]*N)
        m1 = round(float(v1.mean()),1) if v1.notna().sum()>0 else 0
        m2 = round(float(v2.mean()),1) if (c2 and v2.notna().sum()>0) else 0
        if m1>0 or m2>0: dias.append({'label':lbl,'top1':m1,'top2':m2})

    # Cambio en consumo
    cambio = []
    for lbl, c1, c2 in DC['sust_cols']:
        if not c2: continue
        v1 = pd.to_numeric(seg[c1], errors='coerce').fillna(0)
        v2 = pd.to_numeric(seg[c2], errors='coerce').fillna(0)
        mask = v1>0; nc = int(mask.sum())
        if nc<2: continue
        s1=v1[mask]; s2=v2[mask]
        n_abs=int((s2==0).sum()); n_dis=int(((s2>0)&(s2<s1)).sum())
        n_sc=int((s2==s1).sum()); n_emp=int((s2>s1).sum())
        p2 = lambda n: round(n/nc*100,1)
        cambio.append({'label':lbl,'n':nc,
                       'abs':p2(n_abs),'dis':p2(n_dis),'sin':p2(n_sc),'emp':p2(n_emp)})

    # Transgresión
    tr_cols1=[c1 for _,c1,_ in DC['tr_sn']]; tr_cols2=[c2 for _,_,c2 in DC['tr_sn']]
    vif_c1, vif_c2 = DC['vif']
    def has_tr(row, sn_cols, vif_col):
        for c in sn_cols:
            if c and _es_positivo(row.get(c,'')): return True
        if vif_col:
            v = pd.to_numeric(row.get(vif_col, np.nan), errors='coerce')
            return not np.isnan(v) and v>0
        return False
    tr1 = seg.apply(lambda r: int(has_tr(r, tr_cols1, vif_c1)), axis=1)
    tr2 = seg.apply(lambda r: int(has_tr(r, tr_cols2, vif_c2)), axis=1)
    pct_tr1 = pct(int(tr1.sum()), N); pct_tr2 = pct(int(tr2.sum()), N)
    tipos_tr = []
    for lbl, c1, c2 in DC['tr_sn']:
        n1 = int(seg[c1].apply(_es_positivo).sum()) if c1 else 0
        n2 = int(seg[c2].apply(_es_positivo).sum()) if c2 else 0
        tipos_tr.append({'label':lbl,'top1':pct(n1,N),'top2':pct(n2,N)})
    if vif_c1:
        vif1_v = pd.to_numeric(seg[vif_c1], errors='coerce')
        vif2_v = pd.to_numeric(seg[vif_c2], errors='coerce') if vif_c2 else pd.Series([np.nan]*N)
        tipos_tr.append({'label':'VIF','top1':pct(int((vif1_v>0).sum()),N),
                         'top2':pct(int((vif2_v>0).sum()),N)})

    # Salud y vivienda
    sal_psi_c1,sal_psi_c2 = DC['sal_psi']; sal_fis_c1,sal_fis_c2 = DC['sal_fis']
    cal_vid_c1,cal_vid_c2 = DC['cal_vid']; viv1_c1,viv1_c2 = DC['viv1']; viv2_c1,viv2_c2 = DC['viv2']
    salud = [
        {'label':'Salud Psicológica','top1':smean(sal_psi_c1),'top2':smean(sal_psi_c2)},
        {'label':'Salud Física',     'top1':smean(sal_fis_c1),'top2':smean(sal_fis_c2)},
        {'label':'Calidad de Vida',  'top1':smean(cal_vid_c1),'top2':smean(cal_vid_c2)},
    ]
    vivienda = [
        {'label':'Lugar estable',      'top1':viv_pct(viv1_c1),'top2':viv_pct(viv1_c2)},
        {'label':'Condiciones básicas','top1':viv_pct(viv2_c1),'top2':viv_pct(viv2_c2)},
    ]

    return dict(N=N, N_total=N_total, seg_tiempo=seg_tiempo, sust_top=sust_top,
                sust=sust, dias=dias, cambio=cambio,
                pct_tr1=pct_tr1, pct_tr2=pct_tr2, tipos_tr=tipos_tr,
                salud=salud, vivienda=vivienda)

# ── Gráficos ──────────────────────────────────────────────────────────────────
def g_torta(d):
    if not d['sust']: return None
    labs=[s['label'] for s in d['sust']]; vals=[s['n'] for s in d['sust']]
    fig,ax=plt.subplots(figsize=(5,4))
    wedges,_,at=ax.pie(vals,labels=None,colors=PIE_COLS[:len(vals)],
        autopct=lambda p:f'{p:.1f}%' if p>3 else '',startangle=140,pctdistance=0.72,
        wedgeprops={'edgecolor':'white','linewidth':2})
    for a in at: a.set_fontsize(9); a.set_color('white'); a.set_fontweight('bold')
    ax.legend(wedges,[f'{l} (n={v})' for l,v in zip(labs,vals)],
              loc='lower center',bbox_to_anchor=(0.5,-0.18),ncol=2,fontsize=8,frameon=False)
    ax.set_aspect('equal'); fig.patch.set_facecolor('white'); fig.tight_layout()
    return fig

def g_dias(d):
    if not d['dias']: return None
    labs=[x['label'] for x in d['dias']]
    t1=[x['top1'] for x in d['dias']]; t2=[x['top2'] for x in d['dias']]
    x=np.arange(len(labs)); ww=0.35
    fig,ax=plt.subplots(figsize=(max(5,len(labs)*0.9),3.5))
    b1=ax.bar(x-ww/2,t1,ww,color=MC_T1,label='Ingreso (TOP1)',zorder=3)
    b2=ax.bar(x+ww/2,t2,ww,color=MC_T2,label='Seguimiento (TOP2)',zorder=3)
    for bar,v in zip(list(b1)+list(b2),t1+t2):
        if v>0: ax.text(bar.get_x()+bar.get_width()/2,bar.get_height()+0.2,
                        f'{v}d',ha='center',va='bottom',fontsize=8,fontweight='bold',color='#333')
    ax.set_xticks(x); ax.set_xticklabels(labs,fontsize=8)
    ax.set_ylabel('Promedio días (0–28)',fontsize=8,color='#595959')
    ax.legend(fontsize=8,frameon=False)
    _ax_style(ax); fig.patch.set_facecolor('white'); fig.tight_layout()
    return fig

def g_cambio(d):
    if not d['cambio']: return None
    labs=[x['label'] for x in d['cambio']]
    abs_=[x['abs'] for x in d['cambio']]; dis_=[x['dis'] for x in d['cambio']]
    sc_=[x['sin'] for x in d['cambio']]; emp_=[x['emp'] for x in d['cambio']]
    x=np.arange(len(labs))
    fig,ax=plt.subplots(figsize=(max(5,len(labs)*0.9),3.5))
    ax.bar(x,abs_,color=MC_ABS,label='Abstinencia',zorder=3)
    ax.bar(x,dis_,bottom=abs_,color=MC_DIS,label='Disminuyó',zorder=3)
    ax.bar(x,sc_,bottom=[a+d for a,d in zip(abs_,dis_)],color=MC_SC,label='Sin cambio',zorder=3)
    ax.bar(x,emp_,bottom=[a+d+s for a,d,s in zip(abs_,dis_,sc_)],color=MC_EMP,label='Empeoró',zorder=3)
    for i,(a,d2,s,e) in enumerate(zip(abs_,dis_,sc_,emp_)):
        y=0
        for val,col in [(a,MC_ABS),(d2,MC_DIS),(s,MC_SC),(e,MC_EMP)]:
            if val>9: ax.text(i,y+val/2,f'{val:.0f}%',ha='center',va='center',
                              fontsize=7.5,color='white',fontweight='bold')
            y+=val
    ax.set_xticks(x); ax.set_xticklabels(labs,fontsize=8)
    ax.set_ylim(0,115); ax.set_ylabel('% consumidores al ingreso',fontsize=8,color='#595959')
    ax.legend(fontsize=7.5,frameon=False,ncol=2,loc='upper right')
    _ax_style(ax); fig.patch.set_facecolor('white'); fig.tight_layout()
    return fig

def g_transgresion(d):
    fig,ax=plt.subplots(figsize=(3.5,3))
    cats=['Ingreso\n(TOP1)','Seguimiento\n(TOP2)']
    vals=[d['pct_tr1'],d['pct_tr2']]
    bars=ax.bar(cats,vals,color=[MC_T1,MC_T2],width=0.5,zorder=3)
    for bar,v in zip(bars,vals):
        ax.text(bar.get_x()+bar.get_width()/2,bar.get_height()+0.8,
                f'{v}%',ha='center',va='bottom',fontsize=12,fontweight='bold',color='#333')
    ax.set_ylim(0,max(vals)*1.4 if max(vals)>0 else 1)
    ax.set_ylabel('% personas',fontsize=8,color='#595959')
    _ax_style(ax); fig.patch.set_facecolor('white'); fig.tight_layout()
    return fig

def g_tipos_tr(d):
    if not d['tipos_tr']: return None
    labs=[t['label'] for t in d['tipos_tr']]
    t1=[t['top1'] for t in d['tipos_tr']]; t2=[t['top2'] for t in d['tipos_tr']]
    x=np.arange(len(labs)); ww=0.35
    fig,ax=plt.subplots(figsize=(max(4,len(labs)*0.9),3))
    b1=ax.bar(x-ww/2,t1,ww,color=MC_T1,zorder=3)
    b2=ax.bar(x+ww/2,t2,ww,color=MC_T2,zorder=3)
    for bar,v in zip(list(b1)+list(b2),t1+t2):
        if v>0: ax.text(bar.get_x()+bar.get_width()/2,bar.get_height()+0.3,
                        f'{v}%',ha='center',va='bottom',fontsize=8,fontweight='bold',color='#333')
    ax.set_xticks(x); ax.set_xticklabels(labs,fontsize=8)
    ax.legend([mpatches.Patch(color=MC_T1),mpatches.Patch(color=MC_T2)],
              ['Ingreso','Seguimiento'],fontsize=8,frameon=False)
    _ax_style(ax); fig.patch.set_facecolor('white'); fig.tight_layout()
    return fig

def g_salud(d):
    if not d['salud']: return None
    labs=[s['label'] for s in d['salud']]
    t1=[s['top1'] for s in d['salud']]; t2=[s['top2'] for s in d['salud']]
    y=np.arange(len(labs)); ww=0.35
    fig,ax=plt.subplots(figsize=(5,2.5))
    b1=ax.barh(y-ww/2,t1,ww,color=MC_T1,label='Ingreso',zorder=3)
    b2=ax.barh(y+ww/2,t2,ww,color=MC_T2,label='Seguimiento',zorder=3)
    for bar,v in zip(list(b1)+list(b2),t1+t2):
        ax.text(bar.get_width()+0.1,bar.get_y()+bar.get_height()/2,
                str(v),va='center',fontsize=8,fontweight='bold',color='#333')
    ax.set_yticks(y); ax.set_yticklabels(labs,fontsize=9)
    ax.set_xlim(0,22); ax.axvline(x=10,color='#BFBFBF',linestyle='--',linewidth=0.8)
    ax.legend(fontsize=8,frameon=False,loc='lower right')
    _ax_style(ax,horiz=True); fig.patch.set_facecolor('white'); fig.tight_layout()
    return fig

def g_vivienda(d):
    if not d['vivienda']: return None
    labs=[v['label'] for v in d['vivienda']]
    t1=[v['top1'] for v in d['vivienda']]; t2=[v['top2'] for v in d['vivienda']]
    y=np.arange(len(labs)); ww=0.35
    fig,ax=plt.subplots(figsize=(4.5,2.5))
    b1=ax.barh(y-ww/2,t1,ww,color=MC_T1,label='Ingreso',zorder=3)
    b2=ax.barh(y+ww/2,t2,ww,color=MC_T2,label='Seguimiento',zorder=3)
    for bar,v in zip(list(b1)+list(b2),t1+t2):
        ax.text(bar.get_width()+1,bar.get_y()+bar.get_height()/2,
                f'{v}%',va='center',fontsize=9,fontweight='bold',color='#333')
    ax.set_yticks(y); ax.set_yticklabels(labs,fontsize=9)
    ax.set_xlim(0,115)
    ax.legend(fontsize=8,frameon=False,loc='lower right')
    _ax_style(ax,horiz=True); fig.patch.set_facecolor('white'); fig.tight_layout()
    return fig

# ── Construcción del PPT ──────────────────────────────────────────────────────
def build_pptx(d):
    prs = Presentation()
    prs.slide_width = SLIDE_W; prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    titulo = f'Seguimiento TOP1 vs TOP2 · {NOMBRE_SERVICIO}'
    pie_txt = f'N seguimiento = {d["N"]}  ·  {NOMBRE_SERVICIO}  ·  {PERIODO}'
    pct_seg = round(d['N']/d['N_total']*100,1) if d['N_total']>0 else 0
    st = d['seg_tiempo']

    # ── SLIDE 1: PORTADA ──────────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    add_rect(sl, 0, 0, 3.8, 5.625, C_DARK)
    add_rect(sl, 3.0, 0, 1.5, 5.625, C_MID)
    add_txt(sl, 'Seguimiento', 0.25, 1.5, 3.2, 0.7, size=22, bold=True, color=C_WHITE)
    add_txt(sl, 'TOP1 vs TOP2', 0.25, 2.3, 3.2, 0.5, size=12, color=C_LIGHT)
    add_txt(sl, NOMBRE_SERVICIO.upper(), 0.25, 3.1, 3.2, 0.6, size=13, bold=True, color=C_WHITE)
    add_txt(sl, PERIODO, 0.25, 3.75, 3.2, 0.4, size=11, color=C_LIGHT)
    add_txt(sl, f'N ingreso: {d["N_total"]}  ·  Con seguimiento: {d["N"]} ({pct_seg}%)',
            0.25, 4.3, 3.2, 0.4, size=10, color=C_LIGHT)
    add_txt(sl, 'Monitoreo de Resultados\nde Tratamiento', 4.5, 2.2, 5.2, 1.0,
            size=16, bold=True, color=C_DARK)
    add_txt(sl, f'Sustancia principal: {d["sust_top"]}', 4.5, 3.3, 5.2, 0.4, size=11, color=C_MID)
    if st['mediana']:
        add_txt(sl, f'Tiempo mediano de seguimiento: {st["mediana"]} meses',
                4.5, 3.75, 5.2, 0.4, size=10, color=C_GRAY)

    # ── SLIDE 2: SUSTANCIA PRINCIPAL ──────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    header(sl, titulo)
    add_txt(sl, f'SUSTANCIA PRINCIPAL AL SEGUIMIENTO  ·  {d["sust_top"]}',
            0.25, 0.82, 9.5, 0.35, size=12, bold=True, color=C_MID)
    fig_t = g_torta(d)
    if fig_t: fig_to_pptx(sl, fig_t, 2.0, 1.2, 6.0, 4.0)
    footer(sl, pie_txt)

    # ── SLIDE 3: DÍAS DE CONSUMO ───────────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    header(sl, titulo)
    add_txt(sl, 'PROMEDIO DE DÍAS DE CONSUMO · Ingreso vs Seguimiento (últimas 4 semanas)',
            0.25, 0.82, 9.5, 0.35, size=11, bold=True, color=C_MID)
    fig_d = g_dias(d)
    if fig_d: fig_to_pptx(sl, fig_d, 0.5, 1.2, 9.0, 4.0)
    footer(sl, pie_txt)

    # ── SLIDE 4: CAMBIO EN CONSUMO ────────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    header(sl, titulo)
    add_txt(sl, 'CAMBIO EN EL CONSUMO POR SUSTANCIA · % de consumidores al ingreso',
            0.25, 0.82, 9.5, 0.35, size=11, bold=True, color=C_MID)
    fig_cb = g_cambio(d)
    if fig_cb: fig_to_pptx(sl, fig_cb, 0.5, 1.2, 9.0, 4.0)
    footer(sl, pie_txt)

    # ── SLIDE 5: TRANSGRESIÓN ─────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    header(sl, titulo)
    div_v(sl, 4.95)
    add_txt(sl, 'TRANSGRESIÓN A LA NORMA SOCIAL\nIngreso vs Seguimiento',
            0.25, 0.82, 4.5, 0.6, size=13, bold=True, color=C_GRAY)
    fig_to_pptx(sl, g_transgresion(d), 0.3, 1.5, 4.4, 3.6)
    add_txt(sl, 'DISTRIBUCIÓN POR TIPO DE TRANSGRESIÓN',
            5.15, 0.82, 4.6, 0.35, size=11, bold=True, color=C_GRAY)
    fig_tt = g_tipos_tr(d)
    if fig_tt: fig_to_pptx(sl, fig_tt, 5.1, 1.5, 4.65, 3.6)
    footer(sl, pie_txt)

    # ── SLIDE 6: SALUD Y VIVIENDA ─────────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    header(sl, titulo)
    div_v(sl, 5.05)
    add_txt(sl, 'AUTOPERCEPCIÓN DEL ESTADO DE SALUD\nY CALIDAD DE VIDA (escala 0–20)',
            0.25, 0.82, 4.7, 0.6, size=11, bold=True, color=C_MID)
    fig_sal = g_salud(d)
    if fig_sal: fig_to_pptx(sl, fig_sal, 0.2, 1.5, 4.7, 3.8)
    add_txt(sl, 'CONDICIONES DE VIVIENDA\n(% con condición Sí · TOP1 vs TOP2)',
            5.3, 0.82, 4.5, 0.6, size=11, bold=True, color=C_MID)
    fig_viv = g_vivienda(d)
    if fig_viv: fig_to_pptx(sl, fig_viv, 5.2, 1.5, 4.5, 3.8)
    footer(sl, pie_txt)

    prs.save(OUTPUT_FILE)
    print(f'  ✓ PPT generado: {OUTPUT_FILE}')

# ══════════════════════════════════════════════════════════════════════════════
if __name__ == '__main__':
    print('='*60)
    print('  PPTX Seguimiento TOP  v2.0  —  Iniciando...')
    print('='*60)
    d = cargar_datos()
    print(f'  N={d["N"]}/{d["N_total"]} | {d["sust_top"]} | {NOMBRE_SERVICIO}')
    build_pptx(d)
    print(f'\n{"="*60}')
    print(f'  ✅  LISTO  →  {OUTPUT_FILE}')
    print(f'{"="*60}')
